import os

from flask import Flask, render_template, request, send_file,jsonify,Blueprint
from werkzeug.utils import secure_filename

from flask_cors import cross_origin, CORS

import speech_recognition as sr
from pydub import AudioSegment
from pydub.utils import make_chunks
from moviepy.editor import *
from pptx import Presentation
from fileinput import filename
from PyPDF2 import PdfFileMerger , PdfFileWriter, PdfFileReader
from zipfile import ZipFile
import tabula
#from tabula import convert_into
from PIL import Image
from pdf2image import convert_from_path
#import poppler

import pdfkit  
import pandas as pd
from moviepy.editor import *
import camelot as cam





import pythoncom
import comtypes.client
from pdf2docx import parse
from io import BytesIO
import subprocess
import os.path

#from split_api import api
#from merge_api import api
#from encrypt_api import api
#from rotate_api import api
#from pdfdoc_api import api

from constants import SECRET_KEY, UPLOAD_FOLDER

app = Flask(__name__)  
CORS(app=app)

app.config['SECRET_KEY'] = SECRET_KEY
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


#api = Blueprint('api', __name__, url_prefix='/api')
#app.register_blueprint(api)

@app.route('/')
@app.route('/home') 
def to_home():
    return render_template('index_1.html')

#@app.route('/merge_pdf/')
#def merge_pdf():
#    return render_template('mergepdf.html')

@app.route('/merge_pdf/', methods = ['GET','POST'])  
@cross_origin()
def merge_pdf():
    file1="static/files/merged/Merged_pdf.pdf"
    if request.method == "POST":
        files = request.files.getlist("file[]")
        print(files)
        merger = PdfFileMerger()
        for file in files:
            file.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(file.filename)))
            merger.append(file)
        merger.write(file1)
    #return jsonify(file1)
        file_url = f"{request.url_root}static/files/merged/Merged_pdf.pdf"
        return jsonify({'file_url': file_url})
    #return send_file('static/files/merged/Merged_pdf.pdf', as_attachment=True)

# @app.route('/split_pdf/')
# def split_pdf():
#     return render_template('split.html')

@app.route('/split_pdf/', methods = ['GET','POST'])  
@cross_origin()
def split_pdf_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        inputpdf = PdfFileReader(f)
        compressed_quality = request.form['inlineOptions']

        if request.form['inlineOptions'] ==  'option1':
            compressed_quality = request.form['inlineOptions']
            if compressed_quality == 'option1':
                name_list=[]
                for i in range(inputpdf.numPages):
                    output = PdfFileWriter()
                    output.addPage(inputpdf.getPage(i))
                    name = "static/splitted_" + str(i) + ".pdf"
                    name_list.append(name)

                    outputStream=  open("static/splitted_" + str(i) + ".pdf", "wb")
                    output.write(outputStream)
                    outputStream.close()

                    # Create a ZipFile Object
                    zipObj2 =  ZipFile('static/splitted_pdfs.zip', 'w')
                    for name in name_list:
                        zipObj2.write(name)
                    zipObj2.close()
                
        elif request.form['inlineOptions'] ==  'option2':    
            start_num = int(request.form['typeNumber'])-1
            page_num = int(request.form['typeNumber1'])
            name_list=[]
            for i in range(start_num,page_num):
                output = PdfFileWriter()
                output.addPage(inputpdf.getPage(i))
                name = "static/splitted_" + str(i) + ".pdf"
                name_list.append(name)

                outputStream=  open("static/splitted_" + str(i) + ".pdf", "wb")
                output.write(outputStream)
                outputStream.close()

                # Create a ZipFile Object
                zipObj2 =  ZipFile('static/splitted_pdfs.zip', 'w')
                for name in name_list:
                    zipObj2.write(name)
                zipObj2.close()
        file_url = f"{request.url_root}static/splitted_pdfs.zip"
        return jsonify({'file_url': file_url})
    #return send_file('static/splitted_pdfs.zip',as_attachment=True)

#@app.route('/encrypt_pdf/')
#def encrypt_pdf():
#    return render_template('encryptpdf.html')

@app.route('/encrypt_pdf/', methods = ['GET','POST'])  
def encrypt_pdf_1():
    if request.method == 'POST':  
        f = request.files['file']
        password = request.form['password']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        pdf = PdfFileReader(f)
        name = f.filename[:-4] + ".pdf"
        name_1 = 'static/files/encrypted/' + name
        pdfwrite = PdfFileWriter()
        for page in range(pdf.getNumPages()):
            pdfwrite.addPage(pdf.getPage(page))
        pdfwrite.encrypt(user_pwd = password , owner_pwd=None, use_128bit=True)
        with open("static/files/encrypted/encryped.pdf", 'wb') as fh:
            pdfwrite.write(fh)

    # return jsonify("static/files/encrypted/encryped.pdf")
    file_url = f"{request.url_root}static/files/encrypted/encryped.pdf"
    return jsonify({'file_url': file_url})
    #return send_file("static/files/encrypted/encryped.pdf",attachment_filename=name,as_attachment=True)


import sys

#@app.route('/rotate_pdf/')
#def rotate_pdf():
#    return render_template('rotatepdf.html')

@app.route('/rotate_pdf/', methods = ['GET','POST'])  
def rotate_pdf():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
           
        pdf_read = PdfFileReader(f, strict = False)
        totalpages = pdf_read.numPages
        pagenum = int(request.form['typeNumber']) - 1
        angle_rotate = int(request.form['typeDegree'])
        pdf_write = PdfFileWriter()
        pdf_write.appendPagesFromReader(pdf_read)
        name = f.filename[:-4] + ".pdf"
        
        # Rotate page 90 degrees to the right
        a = pdf_read.getPage(pagenum).rotateClockwise(angle_rotate)
        with open("static/files/rotated/rotated.pdf", 'wb') as fh:
            pdf_write.write(fh)
            if pdf_write.getNumPages():
                # newpath = inFile[:-4] + "_numbered.pdf"
                with open("static/files/rotated/rotated_1.pdf", 'wb') as f_1:
                    pdf_write.write(f_1)

    #return jsonify("static/files/rotated/rotated_1.pdf")
    file_url = f"{request.url_root}static/files/rotated/rotated_1.pdf"
    return jsonify({'file_url': file_url})
    #return send_file("static/files/rotated/rotated_1.pdf", as_attachment=True)

# @app.route('/pdf_doc/')
# def pdf_doc():
#     return render_template('pdfdoc.html')

@app.route('/pdf_doc/', methods = ['GET','POST'])  
@cross_origin()
def pdf_doc():
    if request.method == 'POST':  
        f = request.files['file']
        filename = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(f.filename))
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        #input_file = app.config['UPLOAD_FOLDER'],filename
        word_file= "static/files/pdfdoc/" + secure_filename(f.filename[:-4]) +".docx"
        parse(filename, word_file,start = 0, end = 10)
    file_url = f"{request.url_root}static/files/pdfdoc/" + secure_filename(f.filename[:-4]) +".docx"
    return jsonify({'file_url': file_url})

# @app.route('/pdf_csv/')
# def pdf_csv():
#     return render_template('coming_soon.html')

@app.route('/pdf_csv/', methods = ['GET','POST'])  
@cross_origin()
def pdf_csv_1():
    if request.method == 'POST':  
        f = request.files['file']
        name = "static/files/pdfcsv files/"+secure_filename(f.filename[:-4]) + ".csv"
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        tabula.convert_into(f, name, output_format="csv", pages='all')
    file_url = f"{request.url_root}static/files/pdfcsv files/"+secure_filename(f.filename[:-4]) + ".csv"
    return jsonify({'file_url': file_url})
# @app.route('/pdf_img/')
# def pdf_img():
#     return render_template('pdfimg.html')

@app.route('/pdf_img/', methods = ['GET','POST'])  
@cross_origin()
def pdf_img_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        filename = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(f.filename))
        # Store Pdf with convert_from_path function
        pages = convert_from_path(filename, poppler_path= "static/poppler-0.67.0/bin")
        # Counter to store images of each page of PDF to image
        image_counter = 0
        name_list = []
        # Iterate through all the pages stored above
        for page in pages:
            name = "static/files/pdfimages/images_" + str(image_counter) + ".jpg"
            name_list.append(name)
            # Save the image of the page in system
            page.save(name, 'JPEG')
            # Increment the counter to update filename
            image_counter = image_counter + 1
            # Create a ZipFile Object
            zipObj2 =  ZipFile('static/files/pdfimages/pdf_imgs.zip', 'w')
            for name in name_list:
                zipObj2.write(name)
            zipObj2.close()
    file_url = f"{request.url_root}static/files/pdfimages/pdf_imgs.zip"
    return jsonify({'file_url': file_url})


# @app.route('/img_pdf/')
# def img_pdf():
#     return render_template('imgpdf.html')

@app.route('/img_pdf/', methods = ['GET','POST']) 
@cross_origin() 
def img_pdf_1():
    if request.method == 'POST':  
        files = request.files.getlist('file[]')
        name_list =[]
        new_list =[]
        for file in files:
            file.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(file.filename)))
            name = file.filename[:-4]+".pdf"
            name_list.append(name)
            image = Image.open(file)
            # converting into chunks using img2pdf
            pdf_bytes = image.convert('RGB')
            new_list.append(pdf_bytes)
        output_path = os.path.join(app.config['UPLOAD_FOLDER']+'/imgpdf/view_1.pdf')
        pdf_bytes.save(output_path,save_all=True,append_images=new_list[:-1])
    file_url = f"{request.url_root}static/files/imgpdf/view_1.pdf"
    return jsonify({'file_url': file_url})

# @app.route('/csv_pdf/')
# def csv_pdf():
#     return render_template('csvpdf.html')

@app.route('/csv_pdf/', methods = ['GET','POST'])  
@cross_origin()
def csv_pdf_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        name = os.path.join(app.config['UPLOAD_FOLDER']+"/csvpdf/") +secure_filename(f.filename[:-4])+".html"
        name_1 = os.path.join(app.config['UPLOAD_FOLDER']+"/csvpdf/") +secure_filename(f.filename[:-4])+".pdf"
        csv_file = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(f.filename))
        CSV = pd.read_csv(csv_file,encoding='unicode_escape')
        CSV.to_html(name)
        # configuring pdfkit to point to our installation of wkhtmltopdf 
        config = pdfkit.configuration(wkhtmltopdf = "static/wkhtmltox/bin/wkhtmltopdf.exe")  
        # converting html file to pdf file  
        pdfkit.from_file(input = name, output_path = name_1, configuration = config) 
    file_url = f"{request.url_root}static/files/csvpdf/" +secure_filename(f.filename[:-4])+".pdf"
    return jsonify({'file_url': file_url})

# @app.route('/ppt_pdf/')
# def ppt_pdf():
#     return render_template('coming_soon.html')

@app.route('/ppt_pdf/', methods = ['GET','POST'])  
@cross_origin()
def ppt_pdf_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        pythoncom.CoInitialize()
        # ... inside the thread function ...
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        # powerpoint.Visible = 1
        ppt_name = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        deck = powerpoint.Presentations.Open(ppt_name)
        output_name = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/pptpdf/",secure_filename(f.filename[:-4])+".pdf")
        deck.SaveAs(output_name,32)
        deck.Close()
        powerpoint.Quit()
    file_url = f"{request.url_root}static/files/pptpdf/"+secure_filename(f.filename[:-4])+".pdf"
    return jsonify({'file_url': file_url})

# @app.route('/doc_pdf/')
# def doc_pdf():
#     return render_template('docpdf.html')

@app.route('/doc_pdf/', methods = ['GET','POST'])  
@cross_origin()
def doc_pdf_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        pythoncom.CoInitialize()
        # ... inside the thread function ...
        word = comtypes.client.CreateObject("Word.Application")
        # powerpoint.Visible = 1
        output_name = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/docpdf/",secure_filename(f.filename[:-5]) + ".pdf")
        word_name = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        doc = word.Documents.Open(word_name)
        doc.SaveAs(output_name,17)
        doc.Close()
        word.Quit()
    file_url = f"{request.url_root}static/files/docpdf/"+secure_filename(f.filename[:-5]) + ".pdf"
    return jsonify({'file_url': file_url})

# @app.route('/jpg_png/')
# def jpg_png():
#     return render_template('jpgpng.html')

@app.route('/jpg_png/', methods = ['GET','POST'])  
@cross_origin()
def jpg_png_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        img = Image.open(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        output_path = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/jpgpng/",secure_filename(f.filename[:-4]) + ".png")
        img.save(output_path)
    file_url = f"{request.url_root}static/files/jpgpng/"+secure_filename(f.filename[:-4]) + ".png"
    return jsonify({'file_url': file_url})

# @app.route('/png_jpg/')
# def png_jpg():
#     return render_template('pngjpg.html')

@app.route('/png_jpg/', methods = ['GET','POST'])  
@cross_origin()
def png_jpg_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        img = Image.open(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))).convert('RGB')
        output_path = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/pngjpg/",secure_filename(f.filename[:-4]) + ".jpg")
        img.save(output_path)
    file_url = f"{request.url_root}static/files/pngjpg/"+secure_filename(f.filename[:-4]) + ".jpg"
    return jsonify({'file_url': file_url})

# @app.route('/img_ppt/')
# def img_ppt():
#     return render_template('imgppt.html')

@app.route('/img_ppt/', methods = ['GET','POST'])  
@cross_origin()
def img_ppt_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        placeholder = slide.placeholders[1]
        picture = placeholder.insert_picture(f)
        output_path = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/imgppt/",secure_filename(f.filename[:-4]) + ".pptx")
        prs.save(output_path)   
    file_url = f"{request.url_root}static/files/imgppt/"+secure_filename(f.filename[:-4]) + ".pptx" 
    return jsonify({'file_url': file_url})

# @app.route('/pdf_ppt/')
# def pdf_ppt():
#     return render_template('pdfppt.html')

@app.route('/pdf_ppt/', methods = ['GET','POST'])  
def pdf_ppt_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        name = secure_filename(f.filename[:-4])
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        pdf_file = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        # Create working folder
        base_name = pdf_file.split(".pdf")[0]
        slideimgs = convert_from_path(pdf_file, 300, fmt='ppm', thread_count=2,poppler_path= "static/poppler-0.67.0/bin")
        for i, slideimg in enumerate(slideimgs):
            if i % 10 == 0:
                print("Saving slide: " + str(i))

            imagefile = BytesIO()
            slideimg.save(imagefile, format='tiff')
            imagedata = imagefile.getvalue()
            imagefile.seek(0)
            width, height = slideimg.size

            # Set slide dimensions
            prs.slide_height = height * 9521
            prs.slide_width = width * 9521

            # Add slide
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9521, height=height * 9521)

        prs.save(base_name + '.pptx')
        return jsonify(base_name + '.pptx')
    #return send_file(base_name + '.pptx' , as_attachment=True)


# @app.route('/ppt_img/')
# def ppt_img():
#     return render_template('coming_soon.html')

@app.route('/ppt_img/', methods = ['GET','POST'])  
@cross_origin()
def ppt_img_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],f.filename))
        pythoncom.CoInitialize()
        ### Spcifying input & output
        input_file = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],f.filename)
        output_file = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/pptimg/" + f.filename[:-5]+".jpg")        
        Application = comtypes.client.CreateObject("Powerpoint.Application")
        Presentation = Application.Presentations.Open(input_file)
        Presentation.Slides[1].Export(output_file, "JPEG")
        Presentation.close()
        Application.Quit()        
    file_url = f"{request.url_root}static/files/pptimg/" + f.filename[:-5]+".jpg"
    return jsonify({'file_url': file_url})

# @app.route('/compress_pdf/')
# def compress_pdf():
#     return render_template('compresspdf.html')

# from subprocess import Popen, STDOUT

@app.route('/compress_pdf/', methods = ['POST','GET'])
@cross_origin()
def compress_pdf_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        input_file_path = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        output_file_path = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/compresspdf/" + secure_filename(f.filename[:-5])+".pdf")       
        compressed_quality = request.form['inlineRadioOptions']
        quality = { 0: '/default', 1: '/prepress', 2: '/printer', 3: '/ebook', 4: '/screen'}
        #quality = { "btnradio0": '/default', "btnradio1": '/prepress', "btnradio2": '/printer', "btnradio3": '/ebook', "btnradio4": '/screen'}      
        if compressed_quality == 'option1':
            subprocess.call(["static/ghostscript/gs.exe", '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4', '-dPDFSETTINGS={}'.format(quality[1]),'-dNOPAUSE', '-dQUIET', '-dBATCH', '-sOutputFile={}'.format(output_file_path), input_file_path])      
        elif compressed_quality == 'option2':
            subprocess.call(["static/ghostscript/gs.exe", '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4', '-dPDFSETTINGS={}'.format(quality[2]),'-dNOPAUSE', '-dQUIET', '-dBATCH', '-sOutputFile={}'.format(output_file_path), input_file_path])
        elif compressed_quality == 'option3':
            subprocess.call(["static/ghostscript/gs.exe", '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4', '-dPDFSETTINGS={}'.format(quality[3]),'-dNOPAUSE', '-dQUIET', '-dBATCH', '-sOutputFile={}'.format(output_file_path), input_file_path])
        elif compressed_quality == 'option4':
            subprocess.call(["static/ghostscript/gs.exe", '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4', '-dPDFSETTINGS={}'.format(quality[4]),'-dNOPAUSE', '-dQUIET', '-dBATCH', '-sOutputFile={}'.format(output_file_path), input_file_path])
        elif compressed_quality == 'option0':
            subprocess.call(["static/ghostscript/gs.exe", '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4', '-dPDFSETTINGS={}'.format(quality[0]),'-dNOPAUSE', '-dQUIET', '-dBATCH', '-sOutputFile={}'.format(output_file_path), input_file_path])
    file_url = f"{request.url_root}static/files/compresspdf/" + secure_filename(f.filename[:-5])+".pdf"
    return jsonify({'file_url': file_url})


# @app.route('/mp4_mp3/')
# def mp4_mp3():
#     return render_template('mp4_mp3.html')

@app.route('/mp4_mp3/', methods = ['GET','POST'])  
@cross_origin()
def mp4_mp3_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        input_file =  os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        output_file = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/mp4mp3/" + secure_filename(f.filename[:-4])+".mp3")
        VideoClip = VideoFileClip(input_file)               ## reading mp4 file
        audioclip= VideoClip.audio
        audioclip.write_audiofile(output_file)               ## writing into mp3
        audioclip.close()
        VideoClip.close()
    file_url = f"{request.url_root}static/files/mp4mp3/" + secure_filename(f.filename[:-4])+".mp3"
    return jsonify({'file_url': file_url})
        
# @app.route('/mp4_wav/')
# def mp4_wav():
#     return render_template('mp4_wav.html')

@app.route('/mp4_wav/', methods = ['GET','POST'])  
@cross_origin()
def mp4_wav_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        input_file = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        wav_file = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/mp4wav/" + secure_filename(f.filename[:-4])+".wav")
        VideoClip = VideoFileClip(input_file)               ## reading mp4 file
        audioclip= VideoClip.audio
        audioclip.write_audiofile(wav_file)
        audioclip.close()
        VideoClip.close()
    file_url = f"{request.url_root}static/files//mp4wav/" + secure_filename(f.filename[:-4])+".wav"
    return jsonify({'file_url': file_url})
    
    
# @app.route('/wav_text/')
# def wav_text():
#     return render_template('coming_soon.html')

@app.route('/wav_text/', methods = ['POST','GET'])
@cross_origin()
def wav_text_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))        
        wav_file =  os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        myaudio_wav = AudioSegment.from_file(file = wav_file,format = "wav")
        ## creating chunks from large files
        chunk_length_ms = 30000                            # pydub calculates in millisec
        chunks = make_chunks(myaudio_wav, chunk_length_ms) # Make chunks of "C" seconds       
        folder_name = "audio-chunks"
        # create a directory to store the audio chunks
        if not os.path.isdir(folder_name):
            os.mkdir(folder_name)
        #Export all of the individual chunks as wav files
        audio_full_text = ""
        for i, chunk in enumerate(chunks):
            chunk_name = os.path.join(folder_name, fr"chunk{0}.wav".format(i))
            # print ("exporting", chunk_name)
            chunk.export(chunk_name, format = "wav")            
            # initialize the recognizer from SpeechRecognizer
            recognizer = sr.Recognizer()
            with sr.AudioFile(chunk_name) as source:
                sound = recognizer.record(source,duration =  30)
                #text = recognizer.recognize_google(sound, language = "auto")
                #print(text,'.')
                # try converting it to text
                try:
                    text = recognizer.recognize_google(sound,language='auto')
                except sr.UnknownValueError as e:
                    None
                else:
                    text = f"{text.capitalize()}. "
                    print(text)

                #audio_full_text = audio_full_text + a_new + "."
                    audio_full_text += text 
        print("/n Full_text := ",audio_full_text)
        with open("static/files/wavtext/Audio_Full_text.txt","w+",encoding = 'utf-8') as f:
            f.write(audio_full_text)
            f.close()
    file_url = f"{request.url_root}static/files/wavtext/Audio_Full_text.txt"
    return jsonify({'file_url': file_url})

@app.route('/img_compression/', methods=['GET', 'POST']) 
@cross_origin()
def img_compression_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        # open the image
        picture = Image.open(f)
        output_file = "static/files/compressed_img/" + secure_filename(f.filename) [:-4]+ "_compressed.jpg"
        # picture.save(output_file , "JPEG", optimize = True, quality = 70)
        
        compressed_quality = request.form['inlineRadioOptions']
        #quality = { "btnradio0": '/default', "btnradio1": '/prepress', "btnradio2": '/printer', "btnradio3": '/ebook', "btnradio4": '/screen'}      
        if compressed_quality == 'option1':
            picture.save(output_file , "PNG", optimize = True, quality = 80)
        
        elif compressed_quality == 'option2':
            picture.save(output_file , "PNG", optimize = True, quality = 50)

        elif compressed_quality == 'option3':
            picture.save(output_file , "PNG", optimize = True, quality = 20)

    file_url = f"{request.url_root}static/files/compressed_img/" + secure_filename(f.filename) [:-4]+ "_compressed.jpg"
    return jsonify({'file_url': file_url})




@app.route('/png_webp/', methods = ['GET','POST'])  
@cross_origin()
def png_webp_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        im=Image.open(f).convert('RGB')
        output_path = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/pngwebp/",secure_filename(f.filename[:-4]) + ".webp")
        im.save(output_path,'webp')
    file_url = f"{request.url_root}static/files/pngwebp/" + secure_filename(f.filename[:-4])+".webp"
    return jsonify({'file_url': file_url})

        #return send_file(output_path, as_attachment=True)


@app.route('/webp_png/', methods = ['GET','POST'])  
@cross_origin()
def webp_png_1():
    if request.method == 'POST':  
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        im=Image.open(f).convert('RGB')
        output_path = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER']+"/webppng/",secure_filename(f.filename[:-4]) + ".png")
        im.save(output_path,'png')

    file_url = f"{request.url_root}static/files/webppng/" + secure_filename(f.filename[:-5])+".png"
    return jsonify({'file_url': file_url})

        #return send_file(output_path, as_attachment=True)

@app.route('/pdf_excel/', methods = ['GET','POST'])  
@cross_origin()
def pdf_excel_1():
    if request.method == 'POST':
        f = request.files['file']
        f.save(os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename)))
        output_file = "static/files/pdfexcel/" + secure_filename(f.filename) [:-4]+ ".xlsx"        
        input_file = os.path.join(os.path.abspath(os.path.dirname(__file__)),app.config['UPLOAD_FOLDER'],secure_filename(f.filename))
        table = cam.read_pdf(input_file, pages = 'all', flavor = 'stream')
        print(table)
        table[0].df.to_excel(output_file)   
    file_url = f"{request.url_root}static/files/pdfexcel/" + secure_filename(f.filename) [:-4]+ ".xlsx"
    return jsonify({'file_url': file_url})

@app.route('/about/')
def about():
    return render_template('about.html')

if __name__ == '__main__':  
    app.run(debug=True)